import argparse
import datetime
import pathlib

from pptx import Presentation
from pptx.slide import Slide


def get_notes(slide: Slide) -> str:
    txt = slide.notes_slide.notes_text_frame.text
    return "\n".join(
        f"{idx}. {line.replace('•', '').strip()}"
        for idx, line in enumerate(txt.split("\n"), start=1)
    )


def get_io(args: argparse.Namespace) -> tuple[pathlib.Path, pathlib.Path]:
    infile = args.infile
    outfile = args.outfile
    if outfile is None:
        now = datetime.datetime.now()
        stem = f"{infile.stem} - {now.isoformat()}"
        outfile = infile.with_stem(stem).with_suffix(".txt")

    return infile, outfile


def main() -> None:
    parser = argparse.ArgumentParser(description="Extract PowerPoint notes")
    parser.add_argument(
        "infile", type=pathlib.Path, help="path to PowerPoint Presentation"
    )
    parser.add_argument(
        "--outfile",
        type=pathlib.Path,
        help="path to write '.txt' file to",
        default=None,
    )

    args = parser.parse_args()
    infile, outfile = get_io(args)

    prs = Presentation(infile)

    count = 1
    export = []
    for slide in prs.slides:
        if slide.has_notes_slide:
            notes = get_notes(slide)
            export.append(f"[SLIDE {count}]\n{notes}")
            count += 1

    with open(outfile, "w") as f:
        f.write("\n\n\n".join(export))


if __name__ == "__main__":
    main()
